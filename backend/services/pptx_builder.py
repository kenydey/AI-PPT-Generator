"""
python-pptx 原生表格/图表/排版引擎 - 自定义模板、ChartData 原生图表、素材插入
"""
import io
from pathlib import Path
from typing import Optional

import httpx
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Emu

from PIL import Image
from backend.api.schemas import PPTOutline, SlideOutline
from backend.services.chart_service import render_dot_to_png


def _fetch_image(url: str) -> bytes:
    with httpx.Client(timeout=15.0) as c:
        r = c.get(url)
        r.raise_for_status()
        return r.content


def _fit_image_to_box(img_bytes: bytes, width_emu: int, height_emu: int) -> bytes:
    """Pillow: object-fit cover 逻辑，缩放裁剪后返回 PNG 字节"""
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    w, h = img.size
    # EMU to approximate pixels (96 DPI: 1 inch = 96 px, 1 inch = 914400 EMU)
    tw = int(width_emu * 96 / 914400)
    th = int(height_emu * 96 / 914400)
    if tw <= 0 or th <= 0:
        return img_bytes
    scale = max(tw / w, th / h)
    nw, nh = int(w * scale), int(h * scale)
    img = img.resize((nw, nh), Image.Resampling.LANCZOS)
    left = (nw - tw) // 2
    top = (nh - th) // 2
    img = img.crop((left, top, left + tw, top + th))
    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()


def build_pptx(outline: PPTOutline, template_path: Optional[str | Path] = None) -> bytes:
    """
    将 PPTOutline 渲染为 .pptx 字节流。
    若提供 template_path 则使用该 .pptx 作为母版；否则使用空白演示文稿。
    """
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    layouts = prs.slide_layouts
    title_layout = layouts[0]
    content_layout = layouts[5] if len(layouts) > 5 else layouts[1]

    for i, slide_out in enumerate(outline.slides):
        layout = title_layout if i == 0 else content_layout
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_out.title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        for b in slide_out.bullets[:8]:
            txBox.text_frame.add_paragraph().text = b

        # 原生图表
        if slide_out.chart_type and slide_out.chart_type != "none" and slide_out.chart_data:
            cd = slide_out.chart_data
            categories = cd.get("categories") or []
            series_list = cd.get("series") or []
            if categories and series_list:
                chart_data = CategoryChartData()
                chart_data.categories = categories
                for s in series_list:
                    name = s.get("name", "Series")
                    vals = s.get("values") or []
                    chart_data.add_series(name, vals)
                x, y, cx, cy = Inches(1), Inches(2.5), Inches(6), Inches(4)
                if slide_out.chart_type == "bar":
                    chart_type = XL_CHART_TYPE.BAR_CLUSTERED
                elif slide_out.chart_type == "line":
                    chart_type = XL_CHART_TYPE.LINE
                elif slide_out.chart_type == "pie":
                    chart_type = XL_CHART_TYPE.PIE
                else:
                    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
                slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)

        # 图片素材
        for j, url in enumerate((slide_out.image_urls or [])[:2]):
            try:
                raw = _fetch_image(url)
                width_emu = int(914400 * 4)
                height_emu = int(914400 * 2.25)
                raw = _fit_image_to_box(raw, width_emu, height_emu)
                x = Inches(7 + j * 2.5)
                slide.shapes.add_picture(io.BytesIO(raw), x, Inches(2.5), width=Inches(2.2), height=Inches(1.8))
            except Exception:
                pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
